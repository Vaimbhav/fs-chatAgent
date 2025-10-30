from __future__ import annotations
from pathlib import Path
from agent_app.config import INDEX_EXTS

def is_supported_file(path: str) -> bool:
    return Path(path).suffix.lower() in INDEX_EXTS

def read_text_str(path: str) -> str:
    p = Path(path)
    ext = p.suffix.lower()

    if ext in {".txt", ".md", ".json", ".xml", ".html", ".htm", ".csv"}:
        try:
            return p.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            return ""

    if ext == ".pdf":
        try:
            from pypdf import PdfReader
            rd = PdfReader(str(p))
            return "\n\n".join([(pg.extract_text() or "") for pg in rd.pages])
        except Exception:
            return ""

    if ext in {".doc", ".docx"}:
        try:
            import docx  # python-docx
            doc = docx.Document(str(p))
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception:
            return ""

    if ext in {".ppt", ".pptx"}:
        try:
            from pptx import Presentation  # python-pptx
            prs = Presentation(str(p))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        out.append(shape.text)
            return "\n".join(out)
        except Exception:
            return ""

    # Fallback: try as text
    try:
        return p.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""

# --- alias expected by tests ---
def read_text_for_path(path) -> str:
    """Compatibility wrapper used by tests."""
    return read_text_str(str(path))

#
# from __future__ import annotations
# from pathlib import Path
# from typing import Optional
# from agent_app.config import INDEX_EXTS
#
# def is_supported_file(path: str) -> bool:
#     return Path(path).suffix.lower() in INDEX_EXTS
#
# def read_text_str(path: str) -> str:
#     """
#     Reads text from multiple formats. Imports heavy libs lazily.
#     """
#     p = Path(path)
#     ext = p.suffix.lower()
#
#     if ext in {".txt", ".md", ".json", ".xml", ".html", ".htm", ".csv"}:
#         try:
#             return p.read_text(encoding="utf-8", errors="ignore")
#         except Exception:
#             return ""
#
#     if ext == ".pdf":
#         try:
#             from pypdf import PdfReader
#             rd = PdfReader(str(p))
#             return "\n\n".join([(pg.extract_text() or "") for pg in rd.pages])
#         except Exception:
#             return ""
#
#     if ext in {".doc", ".docx"}:
#         try:
#             import docx  # python-docx
#             doc = docx.Document(str(p))
#             return "\n".join([para.text for para in doc.paragraphs])
#         except Exception:
#             return ""
#
#     if ext in {".ppt", ".pptx"}:
#         try:
#             from pptx import Presentation  # python-pptx
#             prs = Presentation(str(p))
#             out = []
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         out.append(shape.text)
#             return "\n".join(out)
#         except Exception:
#             return ""
#
#     # Unknown type → try as text
#     try:
#         return p.read_text(encoding="utf-8", errors="ignore")
#     except Exception:
#         return ""
#
#
#
#
#
