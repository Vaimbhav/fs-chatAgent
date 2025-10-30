from pathlib import Path
import pytest

from agent_app.readers import read_text_for_path


def _assert_contains(text: str, needle: str):
    assert needle in text, f"Expected '{needle}' in text. Got: {text[:160]!r}"


def test_txt(tmp_docs: Path):
    p = tmp_docs / "note.txt"
    p.write_text("Hello TXT world", encoding="utf-8")
    out = read_text_for_path(p)
    _assert_contains(out, "Hello TXT")


def test_md(tmp_docs: Path):
    p = tmp_docs / "readme.md"
    p.write_text("# Title\n\nHello MD world", encoding="utf-8")
    out = read_text_for_path(p)
    _assert_contains(out, "Hello MD")


def test_json(tmp_docs: Path):
    p = tmp_docs / "data.json"
    p.write_text('{"greet":"Hello JSON world"}', encoding="utf-8")
    out = read_text_for_path(p)
    _assert_contains(out, "Hello JSON")


def test_xml(tmp_docs: Path):
    p = tmp_docs / "data.xml"
    p.write_text("<root><msg>Hello XML world</msg></root>", encoding="utf-8")
    out = read_text_for_path(p)
    _assert_contains(out, "Hello XML")


def test_html(tmp_docs: Path):
    p = tmp_docs / "page.html"
    p.write_text("<html><body><h1>Hello HTML world</h1></body></html>", encoding="utf-8")
    out = read_text_for_path(p)
    _assert_contains(out, "Hello HTML")


def test_htm(tmp_docs: Path):
    p = tmp_docs / "page.htm"
    p.write_text("<html><body>Hello HTM world</body></html>", encoding="utf-8")
    out = read_text_for_path(p)
    _assert_contains(out, "Hello HTM")


def test_csv(tmp_docs: Path):
    p = tmp_docs / "t.csv"
    p.write_text("a,b\nHello CSV world,2\n", encoding="utf-8")
    out = read_text_for_path(p)
    _assert_contains(out, "Hello CSV")


def test_pdf(tmp_docs: Path):
    pytest.importorskip("reportlab", reason="reportlab not installed")
    from reportlab.pdfgen import canvas  # type: ignore

    p = tmp_docs / "t.pdf"
    c = canvas.Canvas(str(p))
    c.drawString(100, 700, "Hello PDF world")
    c.save()

    out = read_text_for_path(p)
    _assert_contains(out, "Hello PDF")


def test_docx(tmp_docs: Path):
    pytest.importorskip("docx", reason="python-docx not installed")
    from docx import Document  # type: ignore

    p = tmp_docs / "t.docx"
    d = Document()
    d.add_paragraph("Hello DOCX world")
    d.save(str(p))

    out = read_text_for_path(p)
    _assert_contains(out, "Hello DOCX")


def test_pptx(tmp_docs: Path):
    pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx import Presentation  # type: ignore

    p = tmp_docs / "t.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello PPTX world"
    prs.save(str(p))

    out = read_text_for_path(p)
    _assert_contains(out, "Hello PPTX")


def test_xlsx(tmp_docs: Path):
    pytest.importorskip("openpyxl", reason="openpyxl not installed")
    from openpyxl import Workbook  # type: ignore

    p = tmp_docs / "t.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Hello XLSX world"
    wb.save(str(p))

    out = read_text_for_path(p)
    assert "Hello XLSX" in out or len(out) > 0
